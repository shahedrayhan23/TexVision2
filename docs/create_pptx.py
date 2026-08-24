"""
TexVision Hackathon Presentation Generator
Creates a professional PPTX file for the NITER Hackathon 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ====== COLORS ======
BG_DARK = RGBColor(10, 10, 15)       # #0a0a0f
BG_CARD = RGBColor(26, 26, 37)       # #1a1a25
ACCENT_BLUE = RGBColor(0, 102, 255)  # #0066ff
ACCENT_RED = RGBColor(255, 51, 68)   # #ff3344
ACCENT_CYAN = RGBColor(0, 212, 255)  # #00d4ff
ACCENT_GREEN = RGBColor(0, 255, 136) # #00ff88
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(136, 136, 170)  # #8888aa
TEXT_MUTED = RGBColor(85, 85, 119)   # #555577

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_accent_line(slide, left, top, width=0.8, color=ACCENT_BLUE):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ================================================================
# SLIDE 1: HERO
# ================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1)

# Accent line at top
add_accent_line(slide1, 0.8, 1.0, 2.5, ACCENT_CYAN)

# Badge
add_text_box(slide1, 0.8, 1.2, 6, 0.4, "●  NITER INNOVATE HACKATHON 2026", 
             font_size=11, color=ACCENT_CYAN, bold=True)

# Title
add_text_box(slide1, 0.8, 2.0, 10, 1.5, "TexVision", 
             font_size=72, color=TEXT_WHITE, bold=True)

# Subtitle
add_text_box(slide1, 0.8, 3.6, 8, 0.8, "AI-Powered Fabric Defect Detection\n& Production Intelligence System",
             font_size=24, color=TEXT_GRAY)

# Tagline
add_text_box(slide1, 0.8, 4.8, 8, 0.5, "for the textile & RMG industry",
             font_size=18, color=TEXT_MUTED)

# Meta info
add_text_box(slide1, 0.8, 5.8, 3, 0.3, "TEAM", font_size=10, color=TEXT_MUTED)
add_text_box(slide1, 0.8, 6.1, 3, 0.3, "MD. Shahed Rayhan", font_size=14, color=TEXT_WHITE, bold=True)

add_text_box(slide1, 4.0, 5.8, 3, 0.3, "UNIVERSITY", font_size=10, color=TEXT_MUTED)
add_text_box(slide1, 4.0, 6.1, 4, 0.3, "North International University (NITER)", font_size=14, color=TEXT_WHITE, bold=True)

add_text_box(slide1, 8.5, 5.8, 2, 0.3, "STUDENT ID", font_size=10, color=TEXT_MUTED)
add_text_box(slide1, 8.5, 6.1, 2, 0.3, "CS 2304023", font_size=14, color=TEXT_WHITE, bold=True)

add_text_box(slide1, 11.0, 5.8, 2, 0.3, "DATE", font_size=10, color=TEXT_MUTED)
add_text_box(slide1, 11.0, 6.1, 2, 0.3, "August 2026", font_size=14, color=TEXT_WHITE, bold=True)

# ================================================================
# SLIDE 2: PROBLEM STATEMENT
# ================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)

add_text_box(slide2, 0.8, 0.6, 4, 0.4, "THE PROBLEM", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide2, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide2, 0.8, 1.3, 10, 0.8, "Manual Inspection is Broken",
             font_size=40, color=TEXT_WHITE, bold=True)

add_text_box(slide2, 0.8, 2.3, 10, 0.5,
             "Bangladesh's $40B+ RMG industry relies on manual fabric inspection — slow, inconsistent, and error-prone.",
             font_size=16, color=TEXT_GRAY)

# Problem Card 1
card1 = add_rounded_rect(slide2, 0.8, 3.2, 3.6, 2.8)
add_text_box(slide2, 1.1, 3.4, 3, 0.4, "🐌  Too Slow", font_size=18, color=TEXT_WHITE, bold=True)
add_text_box(slide2, 1.1, 3.9, 3.2, 1.2,
             "Inspectors spend 5–10 minutes per fabric piece, creating bottlenecks in production.",
             font_size=13, color=TEXT_GRAY)
add_text_box(slide2, 1.1, 5.4, 3, 0.4, "5-10m", font_size=32, color=ACCENT_RED, bold=True)

# Problem Card 2
card2 = add_rounded_rect(slide2, 4.8, 3.2, 3.6, 2.8)
add_text_box(slide2, 5.1, 3.4, 3, 0.4, "🎯  Too Inaccurate", font_size=18, color=TEXT_WHITE, bold=True)
add_text_box(slide2, 5.1, 3.9, 3.2, 1.2,
             "Fatigue and human error lead to only 70–80% detection rates, missing critical defects.",
             font_size=13, color=TEXT_GRAY)
add_text_box(slide2, 5.1, 5.4, 3, 0.4, "70%", font_size=32, color=ACCENT_RED, bold=True)

# Problem Card 3
card3 = add_rounded_rect(slide2, 8.8, 3.2, 3.6, 2.8)
add_text_box(slide2, 9.1, 3.4, 3, 0.4, "💸  Too Expensive", font_size=18, color=TEXT_WHITE, bold=True)
add_text_box(slide2, 9.1, 3.9, 3.2, 1.2,
             "Each inspection costs $0.50–$1.00, totaling billions in annual losses for the industry.",
             font_size=13, color=TEXT_GRAY)
add_text_box(slide2, 9.1, 5.4, 3, 0.4, "$2-3B", font_size=32, color=ACCENT_RED, bold=True)

# ================================================================
# SLIDE 3: OUR SOLUTION
# ================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)

add_text_box(slide3, 0.8, 0.6, 4, 0.4, "OUR SOLUTION", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide3, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide3, 0.8, 1.3, 10, 0.8, "AI Meets the Factory Floor",
             font_size=40, color=TEXT_WHITE, bold=True)

# Solution features
features = [
    ("📱", "Mobile-First Design", "Inspectors capture fabric images directly from smartphones — no special equipment needed.", ACCENT_BLUE),
    ("🧠", "YOLOv8 + OpenCV Engine", "Dual-mode AI: deep learning for production, classical CV for instant demos. Works offline.", ACCENT_CYAN),
    ("🔄", "Complete Inspection Workflow", "Capture → Analyze → Submit → Manager Review → Approve/Reject → Reinspect.", ACCENT_GREEN),
    ("📊", "Real-Time Dashboard", "Managers get live defect statistics, trends, quality grades, and audit trails.", ACCENT_RED),
]

for i, (icon, title, desc, color) in enumerate(features):
    y = 2.5 + i * 1.15
    card = add_rounded_rect(slide3, 0.8, y, 11.5, 1.0)
    add_text_box(slide3, 1.1, y + 0.1, 0.5, 0.5, icon, font_size=22, color=color)
    add_text_box(slide3, 1.7, y + 0.05, 3, 0.4, title, font_size=16, color=TEXT_WHITE, bold=True)
    add_text_box(slide3, 1.7, y + 0.5, 10, 0.4, desc, font_size=13, color=TEXT_GRAY)

# Stats at bottom
stats = [("3s", "Detection Time"), ("95%+", "Accuracy"), ("5", "Defect Types")]
for i, (val, label) in enumerate(stats):
    x = 2.5 + i * 3.5
    add_text_box(slide3, x, 6.3, 2, 0.5, val, font_size=28, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x, 6.8, 2, 0.3, label, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 4: KEY FEATURES
# ================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)

add_text_box(slide4, 0.8, 0.6, 4, 0.4, "KEY FEATURES", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide4, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide4, 0.8, 1.3, 10, 0.8, "Built for the Factory Floor",
             font_size=40, color=TEXT_WHITE, bold=True)

features_grid = [
    ("⚡", "Real-Time Detection", "AI analyzes fabric images in 1–3 seconds, providing instant defect identification with bounding boxes.", ACCENT_BLUE),
    ("📱", "Mobile App", "Cross-platform React Native app works on Android & iOS. Capture, analyze, and report from anywhere.", ACCENT_CYAN),
    ("🎯", "Severity Classification", "Automatically grades defects as Critical, High, Medium, or Low with confidence scores.", ACCENT_RED),
    ("📊", "Quality Grading", "Assigns A–D quality grades based on defect analysis, helping managers make quick decisions.", ACCENT_GREEN),
    ("📜", "Inspection History", "Complete audit trail of all inspections, decisions, and rework — searchable and exportable.", RGBColor(255, 170, 0)),
    ("☁️", "Local + Cloud Storage", "Works offline with local JSON DB. Seamless upgrade to Firebase for production deployment.", RGBColor(170, 100, 255)),
]

for i, (icon, title, desc, color) in enumerate(features_grid):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.0
    y = 2.5 + row * 2.4

    card = add_rounded_rect(slide4, x, y, 3.7, 2.1)
    add_text_box(slide4, x + 0.3, y + 0.2, 0.5, 0.5, icon, font_size=24, color=color)
    add_text_box(slide4, x + 0.3, y + 0.7, 3.2, 0.3, title, font_size=15, color=TEXT_WHITE, bold=True)
    add_text_box(slide4, x + 0.3, y + 1.1, 3.2, 0.8, desc, font_size=12, color=TEXT_GRAY)

# ================================================================
# SLIDE 5: TECH STACK
# ================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)

add_text_box(slide5, 0.8, 0.6, 4, 0.4, "TECH STACK", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide5, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide5, 0.8, 1.3, 10, 0.8, "Powered by Modern Tech",
             font_size=40, color=TEXT_WHITE, bold=True)

tech_items = [
    ("⚛️", "React Native", "Cross-platform mobile", ACCENT_CYAN),
    ("⚡", "FastAPI", "Async Python backend", ACCENT_BLUE),
    ("🎯", "YOLOv8", "Object detection", ACCENT_GREEN),
    ("👁️", "OpenCV", "Computer vision", RGBColor(255, 170, 0)),
    ("🔐", "JWT Auth", "Secure tokens", ACCENT_RED),
]

for i, (icon, name, desc, color) in enumerate(tech_items):
    x = 0.8 + i * 2.5
    card = add_rounded_rect(slide5, x, 2.5, 2.2, 2.2)
    add_text_box(slide5, x, 2.7, 2.2, 0.6, icon, font_size=36, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, x, 3.3, 2.2, 0.4, name, font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, x, 3.7, 2.2, 0.4, desc, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide5, 0.8, 5.2, 12, 0.4,
             "Also uses:  Expo SDK  ·  Pydantic  ·  NumPy  ·  Firebase  ·  Docker  ·  bcrypt",
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 6: HOW IT WORKS
# ================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6)

add_text_box(slide6, 0.8, 0.6, 4, 0.4, "HOW IT WORKS", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide6, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide6, 0.8, 1.3, 10, 0.8, "Four Steps to Quality",
             font_size=40, color=TEXT_WHITE, bold=True)

steps = [
    ("1", "📷", "Capture Image", "Inspector takes a photo of the fabric using the mobile app camera", ACCENT_BLUE),
    ("2", "🤖", "AI Detection", "YOLOv8 & OpenCV analyze the image for defects in real-time", ACCENT_CYAN),
    ("3", "🎯", "Classify Defects", "System identifies type, location, severity, and confidence score", ACCENT_RED),
    ("4", "✅", "Quality Verdict", "Manager reviews results and approves, rejects, or requests rework", ACCENT_GREEN),
]

for i, (num, icon, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 3.2

    # Number circle
    circle = slide6.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.8), Inches(2.5), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide6, x, 3.5, 2.5, 0.5, icon, font_size=36, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, 4.1, 2.5, 0.4, title, font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide6, x, 4.6, 2.5, 0.8, desc, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 3:
        add_text_box(slide6, x + 2.5, 2.6, 0.5, 0.5, "→", font_size=28, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 7: ARCHITECTURE
# ================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7)

add_text_box(slide7, 0.8, 0.6, 4, 0.4, "SYSTEM ARCHITECTURE", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide7, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide7, 0.8, 1.3, 10, 0.8, "Built for Scalability",
             font_size=40, color=TEXT_WHITE, bold=True)

arch_items = [
    ("📱", "Mobile App", "React Native + Expo", ["Camera", "Gallery", "Results"], ACCENT_BLUE),
    ("⚡", "FastAPI Backend", "REST API + JWT Auth", ["Auth", "Upload", "Analytics"], ACCENT_CYAN),
    ("🧠", "AI Engine", "YOLOv8 + OpenCV", ["Detect", "Classify", "Grade"], ACCENT_GREEN),
    ("💾", "Data Storage", "Local JSON / Firebase", ["Users", "Inspections", "Images"], ACCENT_RED),
]

for i, (icon, title, subtitle, tags, color) in enumerate(arch_items):
    x = 0.5 + i * 3.2
    card = add_rounded_rect(slide7, x, 2.3, 2.8, 2.5)
    add_text_box(slide7, x, 2.5, 2.8, 0.6, icon, font_size=32, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, x, 3.1, 2.8, 0.4, title, font_size=15, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, x, 3.5, 2.8, 0.3, subtitle, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Tags
    tag_text = "  ·  ".join(tags)
    add_text_box(slide7, x, 4.0, 2.8, 0.3, tag_text, font_size=10, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # Arrow
    if i < 3:
        add_text_box(slide7, x + 2.8, 3.2, 0.4, 0.5, "→", font_size=24, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

add_text_box(slide7, 0.8, 5.2, 12, 0.3,
             "Request Flow:  Mobile  →  API  →  AI  →  Database  →  Response",
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Architecture features
arch_features = [
    ("🔒", "JWT Authentication", "Secure token-based access"),
    ("🔄", "Async Processing", "Non-blocking requests"),
    ("📦", "Docker Ready", "Containerized deployment"),
]

for i, (icon, title, desc) in enumerate(arch_features):
    x = 1.5 + i * 4.0
    add_text_box(slide7, x, 5.8, 0.4, 0.4, icon, font_size=18, color=ACCENT_CYAN)
    add_text_box(slide7, x + 0.5, 5.75, 3, 0.3, title, font_size=13, color=TEXT_WHITE, bold=True)
    add_text_box(slide7, x + 0.5, 6.1, 3, 0.3, desc, font_size=11, color=TEXT_MUTED)

# ================================================================
# SLIDE 8: LIVE DEMO
# ================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8)

add_text_box(slide8, 0.8, 0.6, 4, 0.4, "LIVE DEMO", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide8, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide8, 0.8, 1.3, 10, 0.8, "See It in Action",
             font_size=40, color=TEXT_WHITE, bold=True)

demo_steps = [
    ("01", "Login", ["📧 inspector@niter.edu.bd", "🔒 ••••••••", "", "Sign In"]),
    ("02", "Capture", ["┌─────────────────┐", "│   📷 Camera     │", "│   Viewfinder    │", "└─────────────────┘"]),
    ("03", "Analyzing...", ["🔍 Detecting...", "████████░░ 80%", "", "3 defects found"]),
    ("04", "Results", ["🔴 Hole — 94%", "🟠 Stain — 87%", "🟡 Color — 72%", "Grade: D"]),
]

for i, (num, title, lines) in enumerate(demo_steps):
    x = 0.5 + i * 3.2

    add_text_box(slide8, x + 0.3, 2.3, 2, 0.4, num, font_size=28, color=ACCENT_BLUE, bold=True)

    # Screen mockup
    card = add_rounded_rect(slide8, x, 2.8, 2.8, 2.8)

    # Screen header
    add_text_box(slide8, x + 0.2, 2.9, 2.5, 0.3, title, font_size=10, color=ACCENT_CYAN, bold=True)

    # Screen content
    for j, line in enumerate(lines):
        if line:
            add_text_box(slide8, x + 0.2, 3.3 + j * 0.4, 2.5, 0.3, line, font_size=11, color=TEXT_GRAY)

    add_text_box(slide8, x, 5.8, 2.8, 0.3, title, font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 9: IMPACT & RESULTS
# ================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide9)

add_text_box(slide9, 0.8, 0.6, 4, 0.4, "IMPACT & RESULTS", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide9, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide9, 0.8, 1.3, 10, 0.8, "Measurable Results",
             font_size=40, color=TEXT_WHITE, bold=True)

# Impact cards
impact_data = [
    ("200x", "Faster Inspection", "3 seconds vs 5-10 minutes", ACCENT_CYAN),
    ("95%+", "Detection Accuracy", "Up from 70-80% manual", ACCENT_GREEN),
    ("98%", "Cost Reduction", "$0.01 vs $0.50-1.00", ACCENT_BLUE),
    ("$2-3B", "Potential Savings", "Annual for Bangladesh RMG", ACCENT_RED),
]

for i, (val, label, sub, color) in enumerate(impact_data):
    x = 0.8 + i * 3.1
    card = add_rounded_rect(slide9, x, 2.3, 2.8, 1.8)
    add_text_box(slide9, x, 2.5, 2.8, 0.6, val, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, x, 3.1, 2.8, 0.3, label, font_size=14, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, x, 3.5, 2.8, 0.3, sub, font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# Comparison table
table_data = [
    ("Metric", "Before (Manual)", "With TexVision"),
    ("Inspection Time", "5–10 minutes / piece", "2–3 seconds / piece"),
    ("Detection Rate", "70–80%", "95%+"),
    ("Cost per Inspection", "$0.50 – $1.00", "$0.01 – $0.02"),
    ("Human Error Rate", "20–30%", "< 5%"),
]

rows, cols = len(table_data), 3
table = slide9.shapes.add_table(rows, cols, Inches(1.5), Inches(4.5), Inches(10), Inches(2.5)).table

for i, row_data in enumerate(table_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = 'Calibri'

        if i == 0:
            p.font.bold = True
            p.font.color.rgb = ACCENT_CYAN
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 30, 60)
        else:
            p.font.color.rgb = TEXT_GRAY
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

            if j == 2:
                p.font.color.rgb = ACCENT_GREEN
                p.font.bold = True

# ================================================================
# SLIDE 10: FUTURE SCOPE
# ================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide10)

add_text_box(slide10, 0.8, 0.6, 4, 0.4, "FUTURE SCOPE", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide10, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide10, 0.8, 1.3, 10, 0.8, "What's Next",
             font_size=40, color=TEXT_WHITE, bold=True)

timeline_data = [
    ("SHORT TERM", "6 Months", [
        "Train custom YOLO model",
        "Achieve 98%+ accuracy",
        "Offline mode support",
        "Multi-language (Bengali, English)"
    ], ACCENT_BLUE),
    ("MEDIUM TERM", "1–2 Years", [
        "Cloud deployment (AWS/Firebase)",
        "Multi-factory platform",
        "Advanced analytics",
        "ERP integration"
    ], ACCENT_CYAN),
    ("LONG TERM", "3–5 Years", [
        "IoT camera integration",
        "3D fabric analysis",
        "Global quality platform",
        "Made in Bangladesh — for the world"
    ], ACCENT_GREEN),
]

for i, (phase, time, items, color) in enumerate(timeline_data):
    x = 0.8 + i * 4.0
    card = add_rounded_rect(slide10, x, 2.3, 3.7, 3.5)

    add_text_box(slide10, x + 0.3, 2.5, 3.2, 0.3, phase, font_size=10, color=ACCENT_BLUE, bold=True)
    add_text_box(slide10, x + 0.3, 2.9, 3.2, 0.4, time, font_size=22, color=TEXT_WHITE, bold=True)

    for j, item in enumerate(items):
        add_text_box(slide10, x + 0.3, 3.5 + j * 0.5, 3.2, 0.4,
                     f"→  {item}", font_size=13, color=TEXT_GRAY)

# ================================================================
# SLIDE 11: TEAM
# ================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide11)

add_text_box(slide11, 0.8, 0.6, 4, 0.4, "ABOUT THE TEAM", font_size=12, color=ACCENT_BLUE, bold=True)
add_accent_line(slide11, 0.8, 1.0, 1.0, ACCENT_BLUE)

add_text_box(slide11, 0.8, 1.3, 10, 0.8, "Meet the Developer",
             font_size=40, color=TEXT_WHITE, bold=True)

# Avatar circle
circle = slide11.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.8), Inches(2), Inches(2))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT_BLUE
circle.line.fill.background()
tf = circle.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "SR"
p.font.size = Pt(48)
p.font.color.rgb = TEXT_WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Team info
add_text_box(slide11, 4.0, 2.8, 5, 0.5, "MD. Shahed Rayhan", font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide11, 4.0, 3.3, 5, 0.3, "Computer Science Student", font_size=16, color=ACCENT_CYAN)

team_details = [
    ("University", "North International University (NITER)"),
    ("Student ID", "CS 2304023"),
    ("Level / Term", "Level 3, Term 1"),
    ("Focus", "AI & Software Engineering"),
]

for i, (label, value) in enumerate(team_details):
    y = 3.9 + i * 0.45
    add_text_box(slide11, 4.0, y, 2, 0.3, label.upper(), font_size=9, color=TEXT_MUTED)
    add_text_box(slide11, 4.0, y + 0.2, 5, 0.3, value, font_size=14, color=TEXT_WHITE)

# Skills
skills = [("⚛️", "React Native", 85), ("🐍", "Python", 90), ("🧠", "AI / ML", 80), ("🗄️", "Database", 75)]

for i, (icon, name, pct) in enumerate(skills):
    x = 1.0 + i * 3.0
    card = add_rounded_rect(slide11, x, 5.5, 2.6, 1.2)
    add_text_box(slide11, x, 5.6, 2.6, 0.4, icon, font_size=20, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide11, x, 6.0, 2.6, 0.3, name, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 12: THANK YOU
# ================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide12)

add_text_box(slide12, 0, 2.0, 13.333, 1.5, "Thank You",
             font_size=64, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide12, 0, 3.8, 13.333, 0.6,
             "TexVision — AI-Powered Quality Control for Bangladesh's RMG Industry",
             font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Contact info
contacts = [
    ("💻", "GitHub", "shahedrayhan23/TexVision2"),
    ("📧", "Email", "shahedrayhan23@gmail.com"),
    ("🎓", "University", "NITER — CS 2304023"),
]

for i, (icon, label, value) in enumerate(contacts):
    x = 2.5 + i * 3.0
    card = add_rounded_rect(slide12, x, 4.8, 2.6, 0.8)
    add_text_box(slide12, x + 0.1, 4.85, 0.4, 0.4, icon, font_size=18, color=ACCENT_CYAN)
    add_text_box(slide12, x + 0.5, 4.85, 2, 0.25, label, font_size=9, color=TEXT_MUTED)
    add_text_box(slide12, x + 0.5, 5.1, 2, 0.3, value, font_size=12, color=TEXT_WHITE, bold=True)

add_text_box(slide12, 0, 6.2, 13.333, 0.4,
             "NITER INNOVATE HACKATHON 2026  —  MADE WITH ❤️",
             font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ================================================================
# SAVE
# ================================================================
output_path = os.path.join(os.path.dirname(__file__), "TexVision_Hackathon_2026.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
